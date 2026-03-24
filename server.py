import os, re, uuid,  asyncio, logging, io
from contextlib import asynccontextmanager
from typing import Optional, Literal, List
from datetime import datetime
from pathlib import Path

import uvicorn, asyncpg
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware

from pydantic import BaseModel, Field
from dotenv import load_dotenv
from langchain_openai import AzureOpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import WebBaseLoader
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings

# Direct parsers — read from BytesIO, no temp files needed
from pypdf import PdfReader
import docx2txt
from pptx import Presentation



from azure.storage.blob import (
    generate_blob_sas,
    BlobSasPermissions,
    BlobServiceClient,
    CorsRule
)
from datetime import timedelta


load_dotenv()

# Logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
logger = logging.getLogger("UniversalDocAPI")
os.environ["USER_AGENT"] = "Universal Document API 2.0"

# ==============================================================================
# CONFIGURATION
# ==============================================================================
class Config:
    """Application configuration from environment variables."""
    AZURE_STORAGE_ACCOUNT = os.getenv("AZURE_STORAGE_ACCOUNT")
    AZURE_STORAGE_KEY = os.getenv("AZURE_STORAGE_KEY")
    AZURE_CONTAINER = os.getenv("AZURE_CONTAINER")
    AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
    AZURE_OPENAI_API_INSTANCE_NAME = os.getenv("AZURE_OPENAI_API_INSTANCE_NAME")
    AZURE_OPENAI_EMBEDDING_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT")
    AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION")
    AZURE_OPENAI_EMBEDDING_MODEL = os.getenv("AZURE_OPENAI_EMBEDDING_MODEL")
    POSTGRES_USER = os.getenv("POSTGRES_USER")
    POSTGRES_PASSWORD = os.getenv("POSTGRES_PASSWORD")
    POSTGRES_HOST = os.getenv("POSTGRES_HOST")
    POSTGRES_PORT = os.getenv("POSTGRES_PORT")
    POSTGRES_DATABASE = os.getenv("POSTGRES_DATABASE")
    CHUNK_SIZE = int(os.getenv("CHUNK_SIZE"))
    CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP"))
    MAX_FILE_SIZE_MB = int(os.getenv("MAX_FILE_SIZE_MB"))
    EMBEDDING_BATCH_SIZE = int(os.getenv("EMBEDDING_BATCH_SIZE"))

config = Config()




class UploadSASRequest(BaseModel):
    college_id: str
    document_id: str
    file_name: str


class ReadSASRequest(BaseModel):
    blob_path: str


# ==============================================================================
# PYDANTIC MODELS
# ==============================================================================
class InsertDocumentRequest(BaseModel):
    document_name: Optional[str] = Field(None, min_length=1, max_length=255)
    document_url: Optional[str] = Field(None, max_length=2000)
    source_type: Literal["base64", "website"]
    source: Literal["Product-AI", "Zox-edu-ai","zox-internal"] = Field("Product-AI", description="Source system for the document")
    document_content: Optional[str] = None
    document_type: Optional[Literal["pdf", "docx", "pptx"]] = None
    chunk_size: Optional[int] = Field(None, ge=100, le=5000)
    chunk_overlap: Optional[int] = Field(None, ge=0, le=1000)


class RetrievalRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=500)
    limit: int = Field(5, ge=1, le=100)
    source: Literal["Product-AI", "Zox-edu-ai"] = Field("Product-AI", description="Source system to filter documents")
    document_type: Optional[str] = None
    job_id: Optional[str] = None


class InsertFromBlobRequest(BaseModel):
    document_id: str = Field(..., description="Dataverse record GUID")
    blob_path: str = Field(..., description="Blob path from /blob/upload-sas response")
    document_name: str = Field(..., min_length=1, max_length=255)
    document_type: Literal["pdf", "docx", "pptx"] = Field(..., description="File type")
    source: Literal["Product-AI", "Zox-edu-ai", "zox-internal"] = Field("Product-AI", description="Source system")
    chunk_size: Optional[int] = Field(None, ge=100, le=5000)
    chunk_overlap: Optional[int] = Field(None, ge=0, le=1000)


# ==============================================================================
# DATABASE CLIENT
# ==============================================================================
class PostgresClient:
    """PostgreSQL connection pool with lazy initialization."""
    
    def __init__(self):
        self._pool: Optional[asyncpg.Pool] = None
        self._lock = asyncio.Lock()

    async def _get_pool(self) -> asyncpg.Pool:
        if self._pool:
            return self._pool
        async with self._lock:
            if self._pool is None:
                logger.info("Initializing PostgreSQL connection pool...")
                self._pool = await asyncpg.create_pool(
                    user=config.POSTGRES_USER, password=config.POSTGRES_PASSWORD,
                    host=config.POSTGRES_HOST, port=config.POSTGRES_PORT,
                    database=config.POSTGRES_DATABASE, min_size=1, max_size=5,
                    command_timeout=60, 
                    server_settings={
                        "application_name": "VectorDBWorkflow",
                        "tcp_keepalives_idle": "60", 
                        "tcp_keepalives_interval": "30",
                        "tcp_keepalives_count": "3",
                    }
                )
                logger.info("✓ PostgreSQL pool initialized")
            return self._pool

    async def _execute(self, method: str, query: str, *args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await getattr(conn, method)(query, *args)

    async def executemany(self, query: str, args): return await self._execute("executemany", query, args)
    async def fetch(self, query: str, *args): return await self._execute("fetch", query, *args)
    async def fetchrow(self, query: str, *args): return await self._execute("fetchrow", query, *args)
    async def fetchval(self, query: str, *args): return await self._execute("fetchval", query, *args)

    async def close(self):
        if self._pool:
            await self._pool.close()
            self._pool = None
            logger.info("✓ PostgreSQL pool closed")

db_pool = PostgresClient()


# ==============================================================================
# EMBEDDINGS & DOCUMENT LOADERS
# ==============================================================================
embeddings_model = AzureOpenAIEmbeddings(
    api_key=config.AZURE_OPENAI_API_KEY,
    azure_endpoint=config.AZURE_OPENAI_API_INSTANCE_NAME,
    azure_deployment=config.AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    api_version=config.AZURE_OPENAI_API_VERSION,
    model=config.AZURE_OPENAI_EMBEDDING_MODEL,
)

# Supported document types for in-memory extraction
SUPPORTED_TYPES = {"pdf", "docx", "pptx"}

# File signature magic bytes for validation
FILE_SIGNATURES = {
    "pdf":  b"%PDF",
    "docx": b"PK\x03\x04",
    "pptx": b"PK\x03\x04",
}


def validate_file_signature(file_bytes: bytes, expected_type: str) -> bool:
    """Validate file by checking magic bytes in the header."""
    sig = FILE_SIGNATURES.get(expected_type.lower())
    if not sig:
        return False
    if len(file_bytes) < len(sig):
        return False
    return file_bytes[:len(sig)] == sig


def download_blob_to_bytes(blob_path: str) -> bytes:
    """Download blob from Azure Storage using account key (server-side, no SAS needed)."""
    blob_name = blob_path.replace(f"{config.AZURE_CONTAINER}/", "", 1)
    connection_string = (
        f"DefaultEndpointsProtocol=https;"
        f"AccountName={config.AZURE_STORAGE_ACCOUNT};"
        f"AccountKey={config.AZURE_STORAGE_KEY};"
        f"EndpointSuffix=core.windows.net"
    )
    blob_service = BlobServiceClient.from_connection_string(connection_string)
    blob_client = blob_service.get_blob_client(
        container=config.AZURE_CONTAINER, blob=blob_name
    )
    return blob_client.download_blob().readall()


# # DEPRECATED — base64 upload path no longer used. All uploads go through /insert-from-blob.
# def extract_text_with_loaders(base64_content: str, document_type: str) -> str:
#     """Extract text from base64-encoded document."""
#     try:
#         content_bytes = base64.b64decode(base64_content, validate=True)
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=f"Invalid base64: {e}")
#
#     max_bytes = config.MAX_FILE_SIZE_MB * 1024 * 1024
#     if len(content_bytes) > max_bytes:
#         raise HTTPException(status_code=413, detail=f"File exceeds {config.MAX_FILE_SIZE_MB}MB limit")
#
#     doc_type = document_type.lower()
#     if doc_type not in DOCUMENT_LOADERS:
#         raise HTTPException(status_code=400, detail=f"Unsupported type: {document_type}")
#
#     loader_class, ext = DOCUMENT_LOADERS[doc_type]
#
#     try:
#         with TemporaryDirectory(dir=config.TEMP_DIR) as temp_dir:
#             temp_path = os.path.join(temp_dir, f"input{ext}")
#             Path(temp_path).write_bytes(content_bytes)
#
#             docs = loader_class(temp_path).load()
#             if not docs:
#                 raise ValueError(f"No content extracted from {doc_type.upper()}")
#
#             text = "\n".join(doc.page_content for doc in docs).strip()
#             if not text:
#                 raise ValueError(f"Empty content from {doc_type.upper()}")
#
#             logger.info(f"Extracted {len(docs)} pages/sections from {doc_type.upper()}")
#             return text
#     except HTTPException:
#         raise
#     except Exception as e:
#         logger.error(f"Extraction error ({document_type}): {e}")
#         raise HTTPException(status_code=400, detail=f"Extraction failed: {e}")



def extract_text_from_bytes(file_bytes: bytes, document_type: str) -> str:
    """Extract text from raw bytes using BytesIO — no temp files, 100% in-memory."""
    max_bytes = config.MAX_FILE_SIZE_MB * 1024 * 1024
    if len(file_bytes) > max_bytes:
        raise HTTPException(status_code=413, detail=f"File exceeds {config.MAX_FILE_SIZE_MB}MB limit")

    doc_type = document_type.lower()
    if doc_type not in SUPPORTED_TYPES:
        raise HTTPException(status_code=400, detail=f"Unsupported type: {document_type}")

    try:
        stream = io.BytesIO(file_bytes)

        if doc_type == "pdf":
            reader = PdfReader(stream)
            text = "\n".join(page.extract_text() or "" for page in reader.pages).strip()
            logger.info(f"Extracted {len(reader.pages)} pages from PDF in memory")

        elif doc_type == "docx":
            text = docx2txt.process(stream).strip()
            logger.info(f"Extracted text from DOCX in memory")

        elif doc_type == "pptx":
            prs = Presentation(stream)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            text = "\n".join(parts).strip()
            logger.info(f"Extracted {len(prs.slides)} slides from PPTX in memory")

        stream.close()

        if not text:
            raise ValueError(f"No text content extracted from {doc_type.upper()}")

        return text
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Extraction error ({document_type}): {e}")
        raise HTTPException(status_code=400, detail=f"Extraction failed: {e}")


async def extract_website_data(url: str) -> str:
    """Extract website content."""
    try:
        docs = await asyncio.to_thread(WebBaseLoader(url).load)
        if not docs:
            raise ValueError("No content loaded from website")

        text = "\n".join(doc.page_content for doc in docs).strip()
        if not text:
            raise ValueError("No text extracted from website")

        logger.info(f"Extracted {len(docs)} page(s) from {url}")
        return text
    except asyncio.TimeoutError:
        raise HTTPException(status_code=408, detail="Scraping timeout")
    except Exception as e:
        logger.error(f"Scraping error ({url}): {e}")
        raise HTTPException(status_code=400, detail=f"Scraping failed: {e}")


# ==============================================================================
# EMBEDDING & DOCUMENT HELPERS
# ==============================================================================
async def embed_documents_batch(embeddings: Embeddings, texts: List[str]) -> List[List[float]]:
    """Embed documents in batches."""
    batch_size = config.EMBEDDING_BATCH_SIZE
    result, total = [], len(texts)
    logger.info(f"Generating embeddings for {total} chunks...")

    for i in range(0, total, batch_size):
        batch_num = (i // batch_size) + 1
        batch = texts[i:i + batch_size]
        try:
            result.extend(await asyncio.to_thread(embeddings.embed_documents, batch))
            logger.info(f"  ✓ Batch {batch_num}/{(total + batch_size - 1) // batch_size}")
        except Exception as e:
            logger.error(f"  ✗ Batch {batch_num} failed: {e}")
            raise HTTPException(status_code=502, detail=f"Embedding error: {e}")
    return result


def create_documents_with_metadata(
    chunks: List[str], doc_name: Optional[str], doc_type: str, doc_url: Optional[str] = None
) -> List[Document]:
    """Create Document objects with metadata."""
    return [
        Document(page_content=chunk, metadata={
            "source": doc_name, "document_type": doc_type, "chunk_id": idx,
            "created_at": datetime.utcnow().isoformat(),
            **({"document_url": doc_url} if doc_url else {})
        })
        for idx, chunk in enumerate(chunks)
    ]


# ==============================================================================
# FASTAPI APPLICATION
# ==============================================================================
async def ensure_azure_cors():
    """Ensure Azure Blob Storage CORS rules are set for Dynamics 365."""
    try:
        if not config.AZURE_STORAGE_ACCOUNT or not config.AZURE_STORAGE_KEY:
            logger.warning("Azure Storage credentials missing; skipping CORS configuration.")
            return

        connection_string = f"DefaultEndpointsProtocol=https;AccountName={config.AZURE_STORAGE_ACCOUNT};AccountKey={config.AZURE_STORAGE_KEY};EndpointSuffix=core.windows.net"
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)

        cors_rule = CorsRule(
            allowed_origins=["https://org94c43c47.crm8.dynamics.com"],
            allowed_methods=["GET", "PUT", "POST", "OPTIONS", "HEAD", "MERGE", "DELETE"],
            allowed_headers=["*"],
            exposed_headers=["*"],
            max_age_in_seconds=3600
        )

        logger.info(f"Ensuring CORS rules for {config.AZURE_STORAGE_ACCOUNT}...")
        blob_service_client.set_service_properties(cors=[cors_rule])
        logger.info("✓ Azure Storage CORS rules validated.")
    except Exception as e:
        logger.error(f"Failed to configure Azure Storage CORS: {e}")


@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("Universal Document API starting...")
    await ensure_azure_cors()
    yield
    await db_pool.close()
    logger.info("Shutdown complete.")

app = FastAPI(
    title="Universal Document API", version="2.0.0",
    description="Modern LangChain implementation with format-specific loaders",
    lifespan=lifespan,
)

# Enable CORS for Dataverse
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://org94c43c47.crm8.dynamics.com"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



# ==============================================================================
# API ROUTES
# ==============================================================================
@app.get("/health", tags=["Health"])
async def health_check():
    """Health check with database connectivity."""
    try:
        result = await db_pool.fetchrow("SELECT NOW() as now")
        return {"status": "healthy", "database": "connected", "timestamp": str(result["now"])}
    except Exception as e:
        raise HTTPException(status_code=503, detail=f"Database error: {e}")


# # DEPRECATED — base64/website insert route. All uploads now go through /insert-from-blob.
# @app.post("/insert", tags=["Documents"])
# async def insert_document(request: InsertDocumentRequest):
#     """Insert document into vector database."""
#     job_id = uuid.uuid4()
#     doc_name = request.document_name or "Unknown"
#
#     try:
#         logger.info(f"Processing job {job_id}: {doc_name} ({request.source_type})")
#
#         # Text Extraction
#         if request.source_type == "base64":
#             if not request.document_name or not request.document_type:
#                 raise HTTPException(status_code=400, detail="Missing name or type for base64")
#             doc_type = request.document_type.lower()
#             text = extract_text_with_loaders(request.document_content, doc_type)
#         elif request.source_type == "website":
#             if not request.document_content or not request.document_content.startswith(("http://", "https://")):
#                 raise HTTPException(status_code=400, detail="Invalid URL")
#             text = await extract_website_data(request.document_content)
#             doc_type, request.document_url = "website", request.document_content
#         else:
#             raise HTTPException(status_code=400, detail="Invalid source_type")
#
#         # Chunking
#         splitter = RecursiveCharacterTextSplitter(
#             chunk_size=request.chunk_size or config.CHUNK_SIZE,
#             chunk_overlap=request.chunk_overlap or config.CHUNK_OVERLAP,
#             add_start_index=True,
#         )
#         chunks = splitter.split_text(text)
#         if not chunks:
#             raise HTTPException(status_code=400, detail="No valid chunks created")
#
#         logger.info(f"Created {len(chunks)} chunks")
#
#         # Embedding & Insert
#         embeddings_list = await embed_documents_batch(embeddings_model, chunks)
#         docs = create_documents_with_metadata(chunks, request.document_name, doc_type, request.document_url)
#
#         await db_pool.executemany(
#             """INSERT INTO education_vector_documents
#                (id, document_name, document_url, embedding, content, job_id, document_type, chunk_index, source, document_id)
#                VALUES ($1, $2, $3, $4::vector, $5, $6, $7::document_type_enum, $8, $9::source_enum, $10)""",
#             [(uuid.uuid4(), d.metadata["source"], request.document_url, str(e), d.page_content,
#               str(job_id), doc_type, d.metadata["chunk_id"], request.source, None) for d, e in zip(docs, embeddings_list)]
#         )
#
#         logger.info(f"✅ Job {job_id}: Inserted {len(chunks)} chunks for '{doc_name}'")
#         return {"success": True, "job_id": str(job_id), "total_chunks": len(chunks), "document_type": doc_type}
#
#     except HTTPException:
#         raise
#     except Exception as e:
#         logger.error(f"❌ Job {job_id} failed: {e}")
#         raise HTTPException(status_code=500, detail=f"Insert failed: {e}")


@app.post("/retrieve", tags=["Documents"])
async def retrieve_documents(request: RetrievalRequest):
    """Retrieve document chunks using vector similarity search."""
    try:
        query_emb = await asyncio.to_thread(embeddings_model.embed_query, request.query)
        query_str = "[" + ",".join(map(str, query_emb)) + "]"

        # Build query with source filter (always applied) and optional filters
        filters, params, idx = ["source = $3::source_enum"], [query_str, request.limit, request.source], 4
        if request.document_type:
            filters.append(f"document_type = ${idx}::document_type_enum")
            params.append(request.document_type)
            idx += 1
        if request.job_id:
            filters.append(f"job_id = ${idx}::uuid")
            params.append(request.job_id)

        where = f"WHERE {' AND '.join(filters)}"
        sql = f"""SELECT id, document_name, document_url, document_type::text, content, job_id, chunk_index, source,
                  1 - (embedding <=> $1::vector)::float as similarity_score
                  FROM education_vector_documents {where} ORDER BY embedding <=> $1::vector LIMIT $2"""

        results = await db_pool.fetch(sql, *params)
        docs = [{**dict(r), "id": str(r["id"]), "job_id": str(r["job_id"])} for r in results]

        logger.info(f"Retrieve: '{request.query[:30]}...' (source: {request.source}) -> {len(docs)} matches")
        return {"success": True, "count": len(docs), "query": request.query, "source": request.source, "documents": docs}

    except Exception as e:
        logger.error(f"Retrieval error: {e}")
        raise HTTPException(status_code=500, detail=f"Retrieval failed: {e}")



@app.post("/insert-from-blob", tags=["Documents"])
async def insert_from_blob(request: InsertFromBlobRequest):
    """Download document from Azure Blob, process it, and upsert into vector database.
    
    If a document with the same document_id exists, all its chunks are deleted
    and replaced with new chunks (transactional upsert).
    """
    job_id = uuid.uuid4()
    # Build document_id to match blob naming: {document_id}_{filename}
    clean_filename = re.sub(r'[^a-zA-Z0-9.-]', '_', request.document_name)
    doc_identifier = f"{request.document_id}_{clean_filename}"

    try:
        logger.info(f"Processing blob job {job_id}: {request.document_name} (doc_id: {doc_identifier})")

        # Step 1: Download blob from Azure Storage
        try:
            blob_bytes = await asyncio.to_thread(download_blob_to_bytes, request.blob_path)
        except Exception as e:
            logger.error(f"Blob download failed: {e}")
            raise HTTPException(status_code=404, detail=f"Blob not found or download failed: {e}")

        # Step 2: Size validation
        max_bytes = config.MAX_FILE_SIZE_MB * 1024 * 1024
        if len(blob_bytes) > max_bytes:
            raise HTTPException(status_code=413, detail=f"File exceeds {config.MAX_FILE_SIZE_MB}MB limit")

        if len(blob_bytes) == 0:
            raise HTTPException(status_code=400, detail="Downloaded file is empty (0 bytes)")

        # Step 3: Magic bytes validation
        if not validate_file_signature(blob_bytes, request.document_type):
            raise HTTPException(
                status_code=400,
                detail=f"File signature does not match expected type '{request.document_type}'. File may be corrupted or misnamed."
            )

        # Step 4: Extract text directly from raw bytes (no base64 needed)
        doc_type = request.document_type.lower()
        text = extract_text_from_bytes(blob_bytes, doc_type)

        # Step 5: Chunk the text
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=request.chunk_size or config.CHUNK_SIZE,
            chunk_overlap=request.chunk_overlap or config.CHUNK_OVERLAP,
            add_start_index=True,
        )
        chunks = splitter.split_text(text)
        if not chunks:
            raise HTTPException(status_code=400, detail="No valid chunks created from document")

        logger.info(f"Created {len(chunks)} chunks from '{request.document_name}'")

        # Step 6: Generate embeddings
        embeddings_list = await embed_documents_batch(embeddings_model, chunks)
        docs = create_documents_with_metadata(chunks, request.document_name, doc_type)

        # Step 7: Transactional upsert — DELETE old chunks + INSERT new chunks
        pool = await db_pool._get_pool()
        async with pool.acquire() as conn:
            async with conn.transaction():
                # Check if document already exists
                existing_count = await conn.fetchval(
                    "SELECT COUNT(*) FROM education_vector_documents WHERE document_id = $1",
                    doc_identifier
                )
                was_update = existing_count > 0

                if was_update:
                    deleted = await conn.execute(
                        "DELETE FROM education_vector_documents WHERE document_id = $1",
                        doc_identifier
                    )
                    logger.info(f"Deleted {deleted} old chunks for document '{doc_identifier}'")

                # Insert new chunks
                await conn.executemany(
                    """INSERT INTO education_vector_documents 
                       (id, document_name, document_url, embedding, content, job_id, 
                        document_type, chunk_index, source, document_id)
                       VALUES ($1, $2, $3, $4::vector, $5, $6, 
                               $7::document_type_enum, $8, $9::source_enum, $10)""",
                    [
                        (
                            uuid.uuid4(), d.metadata["source"], None, str(e), d.page_content,
                            str(job_id), doc_type, d.metadata["chunk_id"], request.source,
                            doc_identifier
                        )
                        for d, e in zip(docs, embeddings_list)
                    ]
                )

        action = "updated" if was_update else "inserted"
        logger.info(f"✅ Job {job_id}: {action.capitalize()} {len(chunks)} chunks for '{request.document_name}' (doc_id: {doc_identifier})")
        return {
            "success": True,
            "job_id": str(job_id),
            "document_id": doc_identifier,
            "total_chunks": len(chunks),
            "document_type": doc_type,
            "was_update": was_update,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Blob insert job {job_id} failed: {e}")
        raise HTTPException(status_code=500, detail=f"Insert from blob failed: {e}")


@app.delete("/document/{document_id}", tags=["Documents"])
async def delete_document(document_id: str):
    """Delete all chunks for a specific document from the vector database."""
    try:
        existing = await db_pool.fetchval(
            "SELECT COUNT(*) FROM education_vector_documents WHERE document_id = $1",
            document_id
        )
        if existing == 0:
            raise HTTPException(status_code=404, detail=f"No document found with id '{document_id}'")

        await db_pool.fetch(
            "DELETE FROM education_vector_documents WHERE document_id = $1",
            document_id
        )
        logger.info(f"🗑️ Deleted {existing} chunks for document '{document_id}'")
        return {"success": True, "document_id": document_id, "deleted_chunks": existing}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Delete failed for '{document_id}': {e}")
        raise HTTPException(status_code=500, detail=f"Delete failed: {e}")


@app.post("/blob/upload-sas", tags=["Blob"])
async def generate_upload_sas(request: UploadSASRequest):
    try:
        # 1. Sanitize the filename sent from the frontend (which is already zx_name + ext)
        # 2. Structure: {college_id}/{document_id}_{clean_filename}
        # This guarantees that the same record (document_id) with the same zx_name always overwrites.
        import re
        clean_filename = re.sub(r'[^a-zA-Z0-9.-]', '_', request.file_name)
        blob_name = f"{request.college_id}/{request.document_id}_{clean_filename}"
        
        # 4. SAS Expiry and Start Time (Add buffer for clock skew)
        now = datetime.utcnow()
        start = now - timedelta(minutes=2)
        expiry = now + timedelta(minutes=15) # Increased to 15 mins for reliability

        sas_token = generate_blob_sas(
            account_name=config.AZURE_STORAGE_ACCOUNT,
            container_name=config.AZURE_CONTAINER,
            blob_name=blob_name,
            account_key=config.AZURE_STORAGE_KEY,
            permission=BlobSasPermissions(write=True, create=True),
            start=start,
            expiry=expiry
        )

        upload_url = (
            f"https://{config.AZURE_STORAGE_ACCOUNT}.blob.core.windows.net/"
            f"{config.AZURE_CONTAINER}/{blob_name}?{sas_token}"
        )

        return {
            "upload_url": upload_url,
            "blob_path": f"{config.AZURE_CONTAINER}/{blob_name}"
        }

    except Exception as e:
        logger.error(f"SAS generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate upload SAS: {str(e)}")

@app.post("/blob/read-sas", tags=["Blob"])
async def generate_read_sas(request: ReadSASRequest):
    try:
        blob_name = request.blob_path.replace(f"{config.AZURE_CONTAINER}/", "")

        expiry = datetime.utcnow() + timedelta(minutes=5)

        sas_token = generate_blob_sas(
            account_name=config.AZURE_STORAGE_ACCOUNT,
            container_name=config.AZURE_CONTAINER,
            blob_name=blob_name,
            account_key=config.AZURE_STORAGE_KEY,
            permission=BlobSasPermissions(read=True),
            expiry=expiry
        )

        read_url = (
            f"https://{config.AZURE_STORAGE_ACCOUNT}.blob.core.windows.net/"
            f"{config.AZURE_CONTAINER}/{blob_name}?{sas_token}"
        )

        return {"read_url": read_url}

    except Exception as e:
        logger.error(f"Read SAS generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate read SAS: {str(e)}")


if __name__ == "__main__":
    uvicorn.run("server:app", host="0.0.0.0", port=8002, reload=True, log_level="info")
